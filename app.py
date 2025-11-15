import os
from google.colab import userdata  # Optional for Colab
import streamlit as st
from PyPDF2 import PdfReader
import docx
from pptx import Presentation
import pandas as pd
from sentence_transformers import SentenceTransformer
import faiss
from fpdf import FPDF
from groq import Groq

# -------------------------
# Streamlit page config
# -------------------------
st.set_page_config(
    page_title="AI Study Assistant",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded"
)

st.markdown("""
<style>
body {
    background: linear-gradient(135deg, #FDE68A, #A5B4FC);
    font-family: 'Comic Sans MS', sans-serif;
}
h1, h2, h3 {
    color: #6B21A8;
}
.stButton>button {
    background-color: #22D3EE;
    color: white;
    font-weight: bold;
}
</style>
""", unsafe_allow_html=True)

st.title("📚 AI Study Assistant")

# -------------------------
# Sidebar
# -------------------------
st.sidebar.header("Settings")
education_level = st.sidebar.selectbox(
    "Select Education Level:",
    ["Primary", "High School", "College", "University (BS/MS)"]
)

uploaded_files = st.file_uploader(
    "Upload study materials (PDF, DOCX, PPTX, XLSX, TXT)", 
    accept_multiple_files=True
)

# -------------------------
# Initialize Groq client
# -------------------------
API_KEY = os.environ.get("studyhelperai")
if not API_KEY:
    st.warning("Groq API key not found in environment variables!")

client = Groq(api_key=API_KEY)

# -------------------------
# Helper functions
# -------------------------
def extract_text(file):
    name = file.name.lower()
    text = ""
    if name.endswith(".pdf"):
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text() + " "
    elif name.endswith(".docx"):
        doc = docx.Document(file)
        for para in doc.paragraphs:
            text += para.text + " "
    elif name.endswith(".pptx"):
        prs = Presentation(file)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + " "
    elif name.endswith(".txt"):
        text = file.read().decode("utf-8")
    elif name.endswith(".xlsx"):
        df = pd.read_excel(file)
        text = " ".join(df.astype(str).values.flatten())
    return text

def generate_summary(text, level="Primary"):
    prompt = f"Explain the following text in a {level}-friendly way and make a short summary:\n{text}"
    response = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.3-70b-versatile",
    )
    return response.choices[0].message.content

def save_pdf(summary_text, filename="summary.pdf"):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(0, 10, summary_text)
    pdf.output(filename)
    return filename

# -------------------------
# Main processing
# -------------------------
if uploaded_files:
    st.success(f"{len(uploaded_files)} file(s) uploaded successfully!")

    all_text = ""
    for file in uploaded_files:
        all_text += extract_text(file) + " "

    st.subheader("Ask questions about your study materials")
    user_question = st.text_input("Enter your question here:")

    if st.button("Get Answer"):
        if user_question:
            # For simplicity, send all text + user question to Groq LLM
            prompt = f"Answer the question based on the following materials, using {education_level} level:\n\n{all_text}\n\nQuestion: {user_question}"
            response = client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama-3.3-70b-versatile",
            )
            answer = response.choices[0].message.content
            st.markdown(f"**Answer:** {answer}")

    if st.button("Generate Summary"):
        summary = generate_summary(all_text, level=education_level)
        st.markdown(f"**Summary:** {summary}")
        pdf_file = save_pdf(summary)
        st.download_button("Download Summary PDF", data=open(pdf_file, "rb"), file_name="summary.pdf")

